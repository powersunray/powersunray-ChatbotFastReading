import os
import pyheif # For .HEIC
import PyPDF2 # For .pdf
import openpyxl # For .xlsx
import subprocess # For .doc
import pytesseract # For image
import trafilatura # For url
from PIL import Image # For image
from database import db
from docx import Document # For .docx
from pptx import Presentation # For .pptx
from langchain_together import TogetherEmbeddings  # Add to create embeddings
from models import DBDocument, Link, DocumentChunk 
from langchain.text_splitter import RecursiveCharacterTextSplitter
from dotenv import load_dotenv
load_dotenv()

# Initialize embeddings model
embeddings_model = TogetherEmbeddings(
    api_key=os.getenv("TOGETHER_AI_API_KEY"),
    model="togethercomputer/m2-bert-80M-32k-retrieval"
)

def extract_text_from_pdf(pdf_path):
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ''
            for page in reader.pages:
                page_text = page.extract_text() or ''
                # Eliminate NUL character and invalid character
                page_text = ''.join(c for c in page_text if c.isprintable())
                text += page_text
            return text
    except Exception as e:
        print(f"Error processing {pdf_path}: {e}")
        return ""

def extract_text_from_url(url):
    try:
        downloaded = trafilatura.fetch_url(url)
        if downloaded:
            result = trafilatura.extract(downloaded)
            # Eliminate NUL character and invalid character
            result = ''.join(c for c in result if c.isprintable())
            return result if result else ""
        else:
            print(f"Failed to download content from {url}")
            return ""
    except Exception as e:
        print(f"Error processing {url}: {e}")
        return ""

def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        # Eliminate NUL character and invalid character
        text = ''.join(c for c in text if c.isprintable())
        return text
    except Exception as e:
        print(f"Error processing {docx_path}: {e}")
        return ""

def extract_text_from_doc(doc_path):
    try:
        result = subprocess.run(['antiword', doc_path], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        text = result.stdout.decode('utf-8')
        if not text.strip():
            raise ValueError("antiword returned empty")
        # Eliminate NUL character and invalid character
        text = ''.join(c for c in text if c.isprintable())
        return text
    except Exception as e:
        print(f"Error processing {doc_path}: {e}")
        return ""
    
def extract_text_from_excel(xlsx_path):
    try:
        workbook = openpyxl.load_workbook(xlsx_path)
        text = ''
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
        # Eliminate NUL character and invalid character
        text = ''.join(c for c in text if c.isprintable())
        return text
    except Exception as e:
        print(f"Error processing {xlsx_path}: {e}")
        return ""

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        # Eliminate NUL character and invalid character
        text = ''.join(c for c in text if c.isprintable())
        return text
    except Exception as e:
        print(f"Error processing {pptx_path}: {e}")
        return ""
    
def extract_text_from_image(image_path):
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img)
        # Eliminate NUL character and invalid character
        text = ''.join(c for c in text if c.isprintable())
        return text
    except Exception as e:
        print(f"Error processing {image_path}: {e}")
        return ""
    
# def convert_heic_to_jpeg(heic_path):
#     heic_file = pyheif.read(heic_path)
#     image = Image.frombytes(
#         heic_file.mode,
#         heic_file.size,
#         heic_file.data,
#         "raw",
#         heic_file.mode,
#         heic_file.stride,
#     )
#     jpeg_path = heic_path.replace('.heic', '.jpg')
#     image.save(jpeg_path, "JPEG")
#     return jpeg_path

def process_and_store_chunks(source, source_type, session_id):
    if source_type == 'file':
        ext = os.path.splitext(source)[1].lower()
        if ext == '.pdf':
            text = extract_text_from_pdf(source)
        elif ext == '.docx':
            text = extract_text_from_docx(source)
        elif ext == '.doc':
            text = extract_text_from_doc(source)
        elif ext == '.xlsx':
            text = extract_text_from_excel(source)
        elif ext == '.pptx':
            text = extract_text_from_pptx(source)
        elif ext in ['.png', '.jpg', '.jpeg']:
            text = extract_text_from_image(source)
        # elif ext == '.heic':
        #     jpeg_path = convert_heic_to_jpeg(source)
        #     text = extract_text_from_image(jpeg_path)
        else:
            print(f"Unsupported file type: {source}")
            return
    elif source_type == 'link':
        text = extract_text_from_url(source)
    else:
        print("Invalid source type")
        return
    
    # Split into chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=150,
    )
    chunks = text_splitter.split_text(text)
    
    # Save chunks into database
    for chunk in chunks:
        # Create embedding for chunk
        embedding = embeddings_model.embed_query(chunk)
        
        chunk_entry = DocumentChunk(
            session_id=session_id,
            document_id=DBDocument.query.filter_by(filepath=source).first().id if source_type == 'file' else None,
            link_id=Link.query.filter_by(url=source).first().id if source_type == 'link' else None,
            chunk_text=chunk,
            embedding=embedding
        )
        db.session.add(chunk_entry)
    db.session.commit()